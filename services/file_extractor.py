"""
file_extractor.py — Extract text from the first page/section of a file.

Supports:
  - PDF   → first page text via pdfplumber
  - DOCX  → first ~500 chars of body text via python-docx
  - PPTX  → text from first slide via python-pptx
  - Other → empty string (fallback to filename + chat context)

Keeps extraction fast and lightweight — only reads as much as needed.
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Max characters to extract — enough context without blowing up the prompt
MAX_EXTRACT_CHARS = 600


def extract_first_page_text(file_path: Path) -> str:
    """
    Extract a short snippet of text from the first page/section of a file.
    Returns an empty string if extraction fails or format is unsupported.
    """
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext == ".docx":
            return _extract_docx(file_path)
        elif ext == ".pptx":
            return _extract_pptx(file_path)
        else:
            return ""
    except Exception as e:
        logger.debug("Could not extract text from %s: %s", file_path.name, e)
        return ""


def _extract_pdf(file_path: Path) -> str:
    """Extract text using PyMuPDF (fitz). If empty, fall back to OCR."""
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF (fitz) not installed — skipping PDF extraction")
        return ""

    try:
        doc = fitz.open(str(file_path))
        if len(doc) == 0:
            doc.close()
            return ""
            
        # Get text from first page
        page = doc[0]
        text = page.get_text("text") or ""
        text = _clean(text)
        
        # If very little text was extracted (e.g., scanned PDF), try OCR
        if len(text) < 30:
            try:
                import pytesseract
                from PIL import Image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text = pytesseract.image_to_string(img).strip()
                ocr_text = _clean(ocr_text)
                if ocr_text and len(ocr_text) > len(text):
                    text = ocr_text
            except Exception as ocr_err:
                logger.debug("OCR fallback failed for %s: %s", file_path.name, ocr_err)
                
        doc.close()
        return text
    except Exception as e:
        logger.debug("PDF extraction failed for %s: %s", file_path.name, e)
        return ""





def _extract_docx(file_path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed — skipping DOCX extraction")
        return ""

    try:
        doc = Document(str(file_path))
        chunks = []
        total = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                chunks.append(text)
                total += len(text)
                if total >= MAX_EXTRACT_CHARS:
                    break
        return _clean("\n".join(chunks))
    except Exception as e:
        logger.debug("DOCX extraction failed for %s: %s", file_path.name, e)
        return ""


def _extract_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping PPTX extraction")
        return ""

    try:
        prs = Presentation(str(file_path))
        if not prs.slides:
            return ""
        # Only first slide
        slide = prs.slides[0]
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text.strip())
        return _clean("\n".join(chunks))
    except Exception as e:
        logger.debug("PPTX extraction failed for %s: %s", file_path.name, e)
        return ""


def _clean(text: str) -> str:
    """Normalize whitespace and truncate to MAX_EXTRACT_CHARS."""
    # Collapse excessive whitespace / newlines
    import re
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = text.strip()
    return text[:MAX_EXTRACT_CHARS]
